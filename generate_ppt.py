from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK=RGBColor(0x0F,0x17,0x2A); ORANGE=RGBColor(0xEA,0x58,0x0C)
BLUE=RGBColor(0x02,0x84,0xC7);  GREEN=RGBColor(0x16,0xA3,0x4A)
PURPLE=RGBColor(0x93,0x33,0xEA);RED=RGBColor(0xDC,0x26,0x26)
YELLOW=RGBColor(0xFB,0xBF,0x24);WHITE=RGBColor(0xFF,0xFF,0xFF)
LGREY=RGBColor(0xF1,0xF5,0xF9); MGREY=RGBColor(0x94,0xA3,0xB8)
LBLUE=RGBColor(0xE0,0xF2,0xFE); LORANGE=RGBColor(0xFF,0xED,0xD5)
LGREEN=RGBColor(0xDC,0xFC,0xE7);LPURPLE=RGBColor(0xF3,0xE8,0xFF)
LRED=RGBColor(0xFE,0xE2,0xE2);  SLATE=RGBColor(0x1E,0x29,0x3B)
TEAL=RGBColor(0x0D,0x94,0x88);  LTEAL=RGBColor(0xCC,0xFB,0xF1)

def blank(bg=RGBColor(0xFA,0xFA,0xF9)):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    b=sl.background.fill; b.solid(); b.fore_color.rgb=bg
    return sl

def R(sl,l,t,w,h,fill):
    s=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.line.fill.background()
    return s

def T(sl,text,l,t,w,h,sz=12,bold=False,color=DARK,align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=True; tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=color
    return tb

def Tmulti(sl,lines,l,t,w,h,sz=11,color=DARK,bold_first=False):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=True; tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(2); p.space_after=Pt(2)
        r=p.add_run(); r.text=line
        r.font.size=Pt(sz); r.font.color.rgb=color
        r.font.bold=(i==0 and bold_first)

def hdr(sl,title,accent=ORANGE):
    R(sl,0,0,13.33,1.05,DARK); R(sl,0,1.05,13.33,0.07,accent)
    T(sl,title,0.3,0.1,12.7,0.85,sz=28,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

def card(sl,l,t,w,h,bg,border,lines,tsize=13,bsize=11):
    R(sl,l,t,w,h,bg); R(sl,l,t,0.07,h,border)
    if lines:
        T(sl,lines[0],l+0.12,t+0.1,w-0.2,0.38,sz=tsize,bold=True,color=border)
        if len(lines)>1:
            Tmulti(sl,lines[1:],l+0.12,t+0.5,w-0.2,h-0.58,sz=bsize,color=DARK)

# ══════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════
sl=blank(DARK)
R(sl,8.8,0,4.53,7.5,ORANGE)
# Diagonal feel
R(sl,8.2,0,0.7,7.5,RGBColor(0x1E,0x29,0x3B))

T(sl,"SAMARPAN",0.4,0.8,8,1.6,sz=78,bold=True,color=WHITE)
T(sl,"Synchronized Multiplayer Learning & Assessment Platform",0.4,2.55,8,0.7,sz=19,color=MGREY)
T(sl,'"Engage. Compete. Learn. — at Scale."',0.4,3.3,7.8,0.55,sz=17,italic=True,color=YELLOW)

# Stats strip
stats=[("$404B","Global EdTech Market"),("60%+","Quiz Drop-off Rate"),
       ("2.4B","Students Worldwide"),("Zero","True Sync Competitors")]
for i,(v,l2) in enumerate(stats):
    y=4.2+i*0.72
    R(sl,0.4,y,3.8,0.62,SLATE)
    T(sl,v, 0.5,y+0.04,1.1,0.54,sz=16,bold=True,color=YELLOW)
    T(sl,l2,1.65,y+0.1,2.5,0.45,sz=11,color=WHITE)

# Right panel
rpanel=[("Hackathon Innovation Track",""),
        ("Team: Aman · Sumit · Anuneet · Janvi",""),
        ("Stage: Working MVP + Beta Users",""),
        ("Seeking: Seed Funding $50,000","")]
T(sl,"Innovation Foundation\nHackathon 2026",8.95,1.0,3.9,1.0,sz=16,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
for i,(line,_) in enumerate(rpanel):
    R(sl,9.0,2.2+i*1.1,3.5,0.9,RGBColor(0x7C,0x2D,0x12))
    T(sl,line,9.1,2.28+i*1.1,3.3,0.72,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 2 – THE PROBLEM
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"THE PROBLEM: WHAT IS MISSING IN TODAY'S PLATFORMS?",ORANGE)

problems=[
    (RED,LRED,"ISOLATION & BOREDOM",
     ["Static platforms (Google Forms, Typeform) provide zero social engagement.",
      "No real-time competition means users treat quizzes as chores, not experiences.",
      "Studies show 67% of learners abandon digital assessments within the first 3 minutes.",
      "Zero motivational feedback loops — no streaks, leaderboards, or live rivalry."]),
    (ORANGE,LORANGE,"BROKEN SYNCHRONIZATION",
     ["Kahoot & Quizizz use client-side timers — easily manipulated and desynced.",
      "Race conditions allow fast-network users to answer before slow-network peers.",
      "No server-enforced state versioning = unfair competitive environments.",
      "Answer reveal during gameplay enables answer-sharing cheating between users."]),
    (PURPLE,LPURPLE,"ZERO RETENTION HOOKS",
     ["No shareable result cards, social challenges, or viral invite mechanics.",
      "No streak tracking, XP systems, or progressive difficulty to retain users.",
      "Hosts receive basic % scores — no per-question drop-off or speed analytics.",
      "No community features: no team play, no rival matchmaking, no history."]),
    (BLUE,LBLUE,"THE MARKET GAP",
     ["Kahoot: engaging but shallow — no anti-cheat, no deep analytics, no B2B API.",
      "Google Forms: powerful but dead — no real-time, no social, no competition.",
      "Enterprise LMS tools: functional but expensive ($500+/mo) and user-hostile.",
      "NO platform bridges high-engagement multiplayer + serious assessment depth."]),
]
for i,(bd,bg,title,bullets) in enumerate(problems):
    col=i%2; row=i//2
    l=0.25+col*6.6; t=1.25+row*2.95
    card(sl,l,t,6.4,2.8,bg,bd,[title]+bullets,tsize=13,bsize=10.5)

R(sl,0,6.88,13.33,0.62,DARK)
T(sl,"The global EdTech market is $404B and growing at 16% YoY. No platform has solved synchronized, fair, engaging assessment at scale. This is our entry point.",
  0.3,6.9,12.7,0.55,sz=11,italic=True,color=YELLOW,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 3 – WHO NEEDS US & MARKET STATS
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"WHO NEEDS SAMARPAN? THE PEOPLE WE ARE BUILT FOR",TEAL)

# 5 user segment cards across top
segments=[
    (BLUE,  LBLUE,  "Students",         "Want learning that feels like a game, not a punishment. Need live rivalry, streaks & instant feedback — not a boring form link."),
    (GREEN, LGREEN, "Teachers",         "Need per-question analytics, not just pass/fail. Want revision sessions that students actually engage with — no app install required."),
    (PURPLE,LPURPLE,"Corporate HR",     "Need training completion + compliance tracking without a $500/mo LMS. Want engaging onboarding quizzes employees actually finish."),
    (ORANGE,LORANGE,"Content Creators", "Run live quiz shows on YouTube/Instagram. Need 200+ players synced in real-time with branded shareable result cards post-session."),
    (RED,   LRED,   "Institutions",     "Schools & colleges need bulk quiz mgmt + portal integration at < Rs 0.50/student/month — not Kahoot's $8/student/year pricing."),
]
for i,(bd,bg,title,body) in enumerate(segments):
    l=0.22+i*2.6
    R(sl,l,1.22,2.48,2.55,bg); R(sl,l,1.22,2.48,0.07,bd)
    T(sl,title,l+0.12,1.27,2.25,0.38,sz=12,bold=True,color=bd)
    T(sl,body, l+0.12,1.68,2.28,2.0, sz=9.5,color=DARK)

# Gap statement
R(sl,0.22,3.9,12.88,0.55,SLATE)
T(sl,'"Every student wants engagement. Every teacher wants insight. Every company wants compliance. No single affordable platform delivers all three — until Samarpan."',
  0.3,3.93,12.7,0.5,sz=11,italic=True,color=YELLOW,align=PP_ALIGN.CENTER)

# Market demand stats — 5 big number tiles
T(sl,"MARKET DEMAND PROOF",0.22,4.55,12.8,0.35,sz=12,bold=True,color=DARK,align=PP_ALIGN.CENTER)

mstats=[
    (BLUE,  "$7.5B",  "India EdTech Market\n(20% YoY Growth)"),
    (GREEN, "250M+",  "Students in India\nLearning Online"),
    (RED,   "67%",    "Drop-off Rate on\nDigital Assessments"),
    (PURPLE,"40%",    "Corporate Training\nNever Completed"),
    (ORANGE,"Rs 150", "Our Target CAC vs\nRs 800+ Industry Avg"),
]
for i,(bc,val,label) in enumerate(mstats):
    x=0.22+i*2.62
    R(sl,x,4.95,2.48,2.3,bc)
    T(sl,val,  x,5.05,2.48,0.9,sz=26,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(sl,label,x,5.95,2.48,1.2,sz=10,color=WHITE,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 4 – SOLUTION
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"THE SOLUTION: PROJECT SAMARPAN — TECHNICAL OVERVIEW",BLUE)

# Architecture flow
R(sl,0.2,1.22,12.9,5.55,LGREY)
T(sl,"SYSTEM ARCHITECTURE & CORE FEATURES",0.3,1.25,12.5,0.4,sz=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)

# 3 architecture boxes
arch=[
    (BLUE,LBLUE,"FRONTEND LAYER","Next.js 14 (App Router)\nReact real-time state mgmt\nTailwind CSS + Framer Motion\nSocket.io Client\nResponsive across all devices"),
    (PURPLE,LPURPLE,"BACKEND ENGINE","Node.js + Express.js\nSocket.io WebSocket Server\nServer-controlled timer engine\nStrict state versioning (v1,v2...)\nRazorpay Payment Gateway"),
    (GREEN,LGREEN,"DATA LAYER","MongoDB Atlas (NoSQL)\nPrisma ORM (type-safe)\nRedis (session & queue cache)\nJWT Auth + bcrypt encryption\nOpenAI API (AI quiz gen)"),
]
for i,(bd,bg,title,body) in enumerate(arch):
    l=0.4+i*4.25
    R(sl,l,1.75,3.95,2.8,bg); R(sl,l,1.75,0.08,2.8,bd)
    T(sl,title,l+0.15,1.8,3.7,0.45,sz=13,bold=True,color=bd)
    T(sl,body, l+0.15,2.3,3.7,2.2,sz=10.5,color=DARK)
    if i<2:
        T(sl,"==>",l+4.05,2.75,0.3,0.5,sz=18,bold=True,color=MGREY,align=PP_ALIGN.CENTER)

# Key innovations
T(sl,"KEY TECHNICAL INNOVATIONS",0.4,4.65,12.5,0.38,sz=13,bold=True,color=DARK)
innov=[
    (RED,   LRED,   "Server-Side Timer",   "All quiz timers run on Node.js server. Clients are pure renderers. Eliminates every form of client-side manipulation."),
    (ORANGE,LORANGE,"State Versioning",    "Every game state carries an incremental version number. Stale updates from lagging clients are auto-rejected."),
    (PURPLE,LPURPLE,"Focus Mode",          "Browser tab-switch detection locks user out of quiz. Prevents answer-searching and screen-sharing cheats."),
    (GREEN, LGREEN, "No-Answer-Reveal",    "Correct answers are never sent to clients during gameplay. Only aggregate stats are pushed post-round only."),
]
for i,(bd,bg,title,body) in enumerate(innov):
    l=0.4+i*3.2
    R(sl,l,5.08,3.0,1.55,bg); R(sl,l,5.08,0.07,1.55,bd)
    T(sl,title,l+0.14,5.1,2.75,0.38,sz=11,bold=True,color=bd)
    T(sl,body, l+0.14,5.5,2.75,1.05,sz=9.5,color=DARK)

# ══════════════════════════════════════════════════
# SLIDE 4 – WHY WE WILL WIN
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"WHY SAMARPAN WILL WIN: COMPETITIVE ADVANTAGE",GREEN)

# Left advantages
advs=[
    (GREEN,LGREEN,"Product-Led Viral Growth Engine",
     ["Every quiz result generates a branded shareable card (LinkedIn, WhatsApp, Instagram).",
      "Each card has a 'Challenge Me' CTA — converting viewers to new users organically.",
      "Referral loop: 1 host quiz -> avg 25 players -> 5 new signups (estimated K-factor: 1.3).",
      "Zero paid acquisition needed in early growth phase — fully organic virality."]),
    (BLUE,LBLUE,"Enterprise-Grade Technical Moat",
     ["WebSocket architecture with horizontal scaling via PM2 cluster mode.",
      "MongoDB handles 50,000+ concurrent document reads/writes during live quizzes.",
      "Graceful reconnection logic: players rejoin mid-quiz without losing state.",
      "99.9% uptime target with Vercel edge + Railway backend infrastructure."]),
    (PURPLE,LPURPLE,"Broad & Defensible Market",
     ["B2C: Students, teachers, content creators, trivia hosts worldwide.",
      "B2B: Schools (50M+ students in India alone), corporate HR, ed-agencies.",
      "Platform-agnostic API in Phase 3 locks in enterprise clients with switching cost.",
      "First-mover advantage in India's synchronized multiplayer EdTech segment."]),
]
for i,(bd,bg,title,bullets) in enumerate(advs):
    card(sl,0.2,1.22+i*2.05,6.5,1.9,bg,bd,[title]+bullets,tsize=12,bsize=10)

# Comparison table
R(sl,6.95,1.22,6.1,5.55,WHITE)
R(sl,6.95,1.22,6.1,0.48,DARK)
T(sl,"COMPETITIVE COMPARISON MATRIX",6.95,1.22,6.1,0.48,sz=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

hdrs=["Capability","Forms","Kahoot","Quizizz","SAMARPAN"]
cws= [2.1,0.85,0.9,0.9,1.1]
cxs=[7.0,9.1,9.95,10.85,11.75]
for ci,h in enumerate(hdrs):
    bg2=DARK if ci==0 else (BLUE if ci==4 else LGREY)
    fc=WHITE if ci in(0,4) else DARK
    R(sl,cxs[ci],1.7,cws[ci]-0.04,0.45,bg2)
    T(sl,h,cxs[ci],1.72,cws[ci],0.41,sz=10,bold=True,color=fc,align=PP_ALIGN.CENTER)

rows=[
    ("Server-side Sync",  "No", "No",  "No",  "YES"),
    ("Anti-Cheat Engine", "No", "No",  "Basic","YES"),
    ("Real-time Teams",   "No", "Yes", "Yes", "YES"),
    ("Deep Host Analytics","Basic","No","Basic","YES"),
    ("Viral Share Cards",  "No", "No",  "No",  "YES"),
    ("SaaS Tier Gating",   "No", "Yes", "Yes", "YES"),
    ("AI Quiz Generation", "No", "No",  "No",  "YES"),
    ("LMS/API Integration","No", "No",  "Paid","YES"),
]
for ri,row in enumerate(rows):
    bg2=LGREY if ri%2==0 else WHITE
    for ci,val in enumerate(row):
        R(sl,cxs[ci],2.15+ri*0.5,cws[ci]-0.04,0.48,bg2)
        vc=GREEN if val=="YES" else (RED if val=="No" else YELLOW)
        T(sl,val,cxs[ci],2.18+ri*0.5,cws[ci],0.44,sz=10,
          color=vc if ci>0 else DARK,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 5 – BUSINESS MODEL
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"BUSINESS MODEL, MONETIZATION & FINANCIAL PROJECTIONS",PURPLE)

# 3 tier cards
tiers=[
    (BLUE,LBLUE,"FREE TIER","Rs 0 / month",
     ["Up to 20 concurrent players","5 quizzes per month limit",
      "Basic accuracy analytics","Standard share cards",
      "Community support only","Purpose: Drive viral acquisition"]),
    (PURPLE,LPURPLE,"PRO TIER","Rs 299 / month",
     ["Up to 200 concurrent players","Unlimited quiz creation",
      "Advanced per-question analytics","AI quiz generation (10/mo)",
      "Priority email support","Custom branding on share cards"]),
    (ORANGE,LORANGE,"ENTERPRISE","Rs 2,499 / month",
     ["Unlimited players & quizzes","Full white-label branding",
      "REST API for LMS integration","Unlimited AI quiz generation",
      "SLA-backed 99.9% uptime","Dedicated account manager"]),
]
for i,(bd,bg,title,price,bullets) in enumerate(tiers):
    l=0.2+i*4.35
    R(sl,l,1.22,4.1,4.0,bg); R(sl,l,1.22,4.1,0.07,bd)
    T(sl,title,l,1.25,4.1,0.48,sz=14,bold=True,color=bd,align=PP_ALIGN.CENTER)
    R(sl,l+0.25,1.73,3.6,0.72,bd)
    T(sl,price,l+0.25,1.77,3.6,0.64,sz=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    Tmulti(sl,["  "+b for b in bullets],l+0.15,2.52,3.8,2.65,sz=10.5,color=DARK)

# MRR Bar Chart
R(sl,0.2,5.35,12.9,1.88,WHITE)
T(sl,"MONTHLY RECURRING REVENUE GROWTH PROJECTION (YEAR 1)",0.4,5.38,12.5,0.38,sz=12,bold=True,color=DARK)

bars=[("Mo 1","$0",0.0,MGREY),("Mo 3","$500",0.033,BLUE),("Mo 5","$1.2K",0.08,BLUE),
      ("Mo 7","$2K",0.133,PURPLE),("Mo 9","$3.2K",0.213,PURPLE),
      ("Mo 11","$4K",0.267,GREEN),("Mo 12","$5K+",0.333,GREEN)]
max_h=1.1
for i,(label,val,frac,bc) in enumerate(bars):
    bx=0.5+i*1.82
    bh=max(max_h*frac,0.05); by=5.35+0.42+max_h-bh
    R(sl,bx,by,1.4,bh,bc)
    T(sl,val,bx,by-0.3,1.4,0.28,sz=10,bold=True,color=bc,align=PP_ALIGN.CENTER)
    T(sl,label,bx,6.95,1.4,0.25,sz=10,color=DARK,align=PP_ALIGN.CENTER)

# Unit economics note
T(sl,"Target: Break-even at 200 Pro users | CAC < Rs 150 | LTV:CAC ratio > 10x | Gross Margin ~82%",
  0.3,7.2,12.7,0.25,sz=10,italic=True,color=MGREY,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 6 – ROADMAP
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"STRATEGIC ROADMAP: FROM MVP TO MARKET LEADER",ORANGE)

phases=[
    (BLUE,  "PHASE 1\nCOMPLETED","Q1-Q2 2026",
     ["Core WebSocket sync engine","Dynamic lobbies & slot system",
      "Anti-cheat (focus mode + no-reveal)","Host analytics dashboard",
      "Razorpay payment integration","Freemium tier launch"]),
    (GREEN, "PHASE 2\nIN PROGRESS","Q3 2026",
     ["Subscription quota enforcement","Pro & Enterprise tier rollout",
      "Beta onboarding: 10 schools","AI quiz gen from PDF upload",
      "Mobile-responsive PWA","100 paying users target"]),
    (PURPLE,"PHASE 3\nNEXT 6 MO","Q4 2026",
     ["OpenAI adaptive difficulty","Streak & XP gamification system",
      "Challenge invite & referral loop","Corporate HR pilot (5 companies)",
      "1,000 MAU milestone","Rs 1L MRR target"]),
    (ORANGE,"PHASE 4\n12 MONTHS","Q1 2027",
     ["Public REST API for LMS vendors","React Native iOS & Android app",
      "Multi-language quiz support","50+ school partnerships",
      "10,000+ MAU milestone","Series A fundraise preparation"]),
]
for i,(color,phase,period,bullets) in enumerate(phases):
    l=0.2+i*3.28
    R(sl,l,1.22,3.1,0.75,color)
    T(sl,phase,l,1.22,3.1,0.75,sz=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(sl,period,l,1.97,3.1,0.3,sz=10,italic=True,color=color,align=PP_ALIGN.CENTER)
    R(sl,l,2.27,3.1,3.55,LGREY); R(sl,l,2.27,0.07,3.55,color)
    Tmulti(sl,["  - "+b for b in bullets],l+0.12,2.32,2.9,3.45,sz=10.5,color=DARK)
    if i<3:
        T(sl,">>",l+3.15,3.5,0.18,0.5,sz=14,bold=True,color=color,align=PP_ALIGN.CENTER)

# KPI strip
R(sl,0,6.0,13.33,1.5,DARK)
T(sl,"TARGET MILESTONES",0.3,6.03,12.7,0.38,sz=13,bold=True,color=YELLOW,align=PP_ALIGN.CENTER)
kpis=[("10,000+","MAU by Year 1"),("50+","Schools in Beta"),
      ("Rs 4L","MRR by Mo 12"),("82%","Gross Margin"),("< Rs 150","Customer CAC")]
for i,(v,l2) in enumerate(kpis):
    x=0.3+i*2.6
    T(sl,v,x,6.45,2.4,0.5,sz=16,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(sl,l2,x,6.95,2.4,0.38,sz=10,color=MGREY,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 7 – TEAM & ASK
# ══════════════════════════════════════════════════
sl=blank()
hdr(sl,"THE TEAM, TRACTION & FUNDRAISING ASK",RED)

team=[
    ("AM","Aman Maurya","Team Lead & Full Stack Engineer",
     "Architected the entire WebSocket sync engine, Prisma schema design, REST APIs, and server-side timer logic. Led all technical decisions.",
     BLUE),
    ("SK","Sumit Kumar","Frontend & UI/UX Engineer",
     "Built the Next.js 14 frontend, real-time lobby UI, Tailwind design system, Framer Motion animations, and responsive mobile layout.",
     PURPLE),
    ("AG","Anuneet Gupta","Backend & Payments Engineer",
     "Implemented Razorpay subscription integration, API quota enforcement, JWT auth flows, and billing webhook handlers.",
     GREEN),
    ("JS","Janvi Sahu","AI & ML Integration Engineer",
     "Integrated OpenAI API for quiz generation from PDFs, built adaptive difficulty scoring logic, and designed the AI prompt pipeline.",
     ORANGE),
]
for i,(init,name,role,desc,color) in enumerate(team):
    t2=1.22+i*1.48
    R(sl,0.2,t2,0.82,0.82,color)
    T(sl,init,0.2,t2,0.82,0.82,sz=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    R(sl,1.1,t2,5.55,1.38,LGREY); R(sl,1.1,t2,0.07,1.38,color)
    T(sl,name,1.22,t2+0.05,3.5,0.38,sz=13,bold=True,color=DARK)
    T(sl,role,1.22,t2+0.43,5.0,0.32,sz=10.5,bold=True,color=color)
    T(sl,desc,1.22,t2+0.75,5.3,0.58,sz=9.5,color=DARK)

# Right: Traction + Ask
R(sl,6.9,1.22,6.15,5.8,DARK)

# Traction
R(sl,7.0,1.32,5.95,0.42,SLATE)
T(sl,"TRACTION & VALIDATION",7.0,1.33,5.95,0.4,sz=12,bold=True,color=YELLOW,align=PP_ALIGN.CENTER)
traction=[
    "Working MVP deployed and live",
    "Beta tested with 3 educational institutions",
    "Core multiplayer engine handles 50+ simultaneous users",
    "Razorpay payment pipeline fully integrated",
    "AI quiz generation (PDF -> quiz) prototype complete",
    "Admin dashboard + analytics portal built",
]
Tmulti(sl,["  + "+t for t in traction],7.05,1.78,5.85,1.95,sz=10.5,color=WHITE)

# Ask
R(sl,7.0,3.8,5.95,0.42,ORANGE)
T(sl,"FUNDRAISING ASK",7.0,3.82,5.95,0.38,sz=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
R(sl,7.15,4.27,5.65,0.82,SLATE)
T(sl,"$50,000  /  Rs 40 Lakhs",7.15,4.3,5.65,0.76,sz=20,bold=True,color=YELLOW,align=PP_ALIGN.CENTER)
T(sl,"SEED ROUND",7.15,5.1,5.65,0.3,sz=11,color=MGREY,align=PP_ALIGN.CENTER)

uses=[("50%","Engineering & Server Scale — WebSocket infra, Redis, CDN, DevOps pipeline"),
      ("30%","B2B Sales & Marketing — School outreach, content, growth hacking"),
      ("20%","Operational Runway — Team sustenance for 12-18 month runway")]
for i,(pct,label) in enumerate(uses):
    R(sl,7.05,5.48+i*0.52,5.85,0.48,SLATE)
    T(sl,pct, 7.1,5.5+i*0.52,0.85,0.44,sz=13,bold=True,color=YELLOW)
    T(sl,label,7.98,5.52+i*0.52,4.85,0.42,sz=9.5,color=WHITE)

T(sl,'"Let\'s make learning something people actually look forward to."',
  6.95,6.72,6.1,0.72,sz=11,italic=True,color=YELLOW,align=PP_ALIGN.CENTER)

out=r"c:\Users\maury\OneDrive\Desktop\sam\Samarpan_Final_v2.pptx"
prs.save(out)
print("Saved:",out)
